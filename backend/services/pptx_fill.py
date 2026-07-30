"""
Single source of truth for placeholder replacement.

Both the preview path and the final-render path call fill_template(), so the
two can never drift apart.

--------------------------------------------------------------------------
Why the JPG text used to overlap (root cause + real fix)
--------------------------------------------------------------------------
PowerPoint templates commonly use "Shrink text on overflow" autofit. When you
type into a placeholder in PowerPoint, PowerPoint re-measures the text on
every keystroke and writes the resulting scale to
<a:normAutofit fontScale="..." lnSpcReduction="..."/> on the shape's bodyPr.
That cached value was computed for the ORIGINAL placeholder text.

LibreOffice does not re-measure anything -- it just trusts whatever scale is
cached in the XML. Two bad things can happen depending on what we do:

  1. Leave the cached scale in place: it was computed for short placeholder
     text, so real (usually longer) content still overflows the box.
  2. Strip the cached scale (the previous "fix" in this file) and just turn
     on word_wrap: LibreOffice now renders at the shape's full, un-shrunk
     font size. word_wrap only prevents horizontal overflow -- it does
     nothing for vertical overflow. Longer text wraps to more lines than
     the box is tall, and those extra lines spill downward on top of
     whatever shape sits below (e.g. a "Title" box printing over the
     "Description" box beneath it). This is the overlap you were seeing.

Neither approach actually shrinks the text. The real fix is to do PowerPoint's
job ourselves: estimate how many lines the new text will wrap to inside the
box, and if that would overflow the box height, scale the font size down
(and re-apply it to the runs directly) before LibreOffice ever renders it.
Baking the final size into the runs -- rather than relying on a fontScale
hint -- also means this renders correctly in PowerPoint, LibreOffice, and any
other consumer, since none of them need to interpret the autofit hint at all.

This is a heuristic (no real font metrics are loaded), the same category of
approximation the rest of this codebase already relies on for layout, but it
is a *closed-loop* shrink -- it directly targets "will this fit" -- instead
of a fixed strip-and-hope.
"""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Average glyph width as a fraction of font size, for typical UI/deck fonts
# (Calibri, Arial, Segoe UI, etc). Slightly wider for bold since bold glyphs
# run heavier. This is intentionally conservative (estimates a bit wide) so
# we shrink slightly more often than strictly necessary rather than risk
# residual overflow.
_AVG_CHAR_WIDTH_FACTOR = 0.52
_AVG_CHAR_WIDTH_FACTOR_BOLD = 0.58

_LINE_HEIGHT_FACTOR = 1.2  # single line spacing, in multiples of font size
_DEFAULT_INSET_PT = 7.2  # PowerPoint's default 0.1in text-box inset
_MIN_FONT_SIZE_PT = 8.0
_MIN_SCALE = 0.4  # never shrink below 40% -- past this it reads as a bug, not a fit
_MAX_SCALE = 1.0


# ── styling ────────────────────────────────────────────────────────────────
def _hex_to_rgb(hex_color):
    try:
        h = str(hex_color).lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return None


def _apply_style(run, style):
    if not style:
        return
    try:
        font = run.font
        if style.get("bold") is not None:
            font.bold = bool(style["bold"])
        if style.get("italic") is not None:
            font.italic = bool(style["italic"])
        if style.get("underline") is not None:
            font.underline = bool(style["underline"])
        if style.get("fontSize"):
            font.size = Pt(float(style["fontSize"]))
        if style.get("color"):
            rgb = _hex_to_rgb(style["color"])
            if rgb:
                font.color.rgb = rgb
        # fontFamily is already the exact PPT font name -- use directly.
        if style.get("fontFamily"):
            font.name = style["fontFamily"]
    except Exception as e:
        print(f"Style apply error: {e}")


# ── shrink-to-fit ────────────────────────────────────────────────────────────
def _box_inset_pt(bodyPr, side_attr):
    """Read a bodyPr inset (lIns/rIns/tIns/bIns) in points, EMU -> pt."""
    if bodyPr is None:
        return _DEFAULT_INSET_PT
    val = bodyPr.get(side_attr)
    if val is None:
        return _DEFAULT_INSET_PT
    try:
        return Emu(int(val)).pt
    except Exception:
        return _DEFAULT_INSET_PT


def _run_font_size_pt(run, paragraph, default=18.0):
    if run.font and run.font.size:
        return run.font.size.pt
    if paragraph.font and paragraph.font.size:
        return paragraph.font.size.pt
    return default


def _estimate_wrapped_lines(text, avail_width_pt, font_size_pt, bold=False):
    if not text:
        return 1
    if avail_width_pt <= 0:
        return 1
    factor = _AVG_CHAR_WIDTH_FACTOR_BOLD if bold else _AVG_CHAR_WIDTH_FACTOR
    char_width_pt = font_size_pt * factor
    chars_per_line = max(1, int(avail_width_pt / char_width_pt))

    total_lines = 0
    for hard_line in text.split("\n"):
        hard_line = hard_line.strip()
        if not hard_line:
            total_lines += 1
            continue
        total_lines += -(-len(hard_line) // chars_per_line)  # ceil div
    return max(1, total_lines)


def _shrink_to_fit(shape, text_frame):
    """
    Estimate whether the current text overflows the shape and, if so, scale
    every run's font size down so it fits. Applied after replacement, so it
    sees the final text. No-ops if the shape has no usable dimensions.
    """
    width_emu = getattr(shape, "width", None)
    height_emu = getattr(shape, "height", None)
    if not width_emu or not height_emu:
        return  # e.g. table cells: python-pptx doesn't expose cell w/h here;
        # LibreOffice grows row height to fit instead of overlapping, so this
        # is safe to skip.

    body = text_frame._txBody
    bodyPr = body.find(f"{{{_A_NS}}}bodyPr")

    l_ins = _box_inset_pt(bodyPr, "lIns")
    r_ins = _box_inset_pt(bodyPr, "rIns")
    t_ins = _box_inset_pt(bodyPr, "tIns")
    b_ins = _box_inset_pt(bodyPr, "bIns")

    avail_width_pt = Emu(width_emu).pt - l_ins - r_ins
    avail_height_pt = Emu(height_emu).pt - t_ins - b_ins
    if avail_width_pt <= 0 or avail_height_pt <= 0:
        return

    total_lines = 0
    max_font_size = 0.0
    any_bold = False
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            continue
        text = paragraph.text
        font_size = max(_run_font_size_pt(r, paragraph) for r in runs)
        bold = any(bool(r.font.bold) for r in runs if r.font)
        max_font_size = max(max_font_size, font_size)
        any_bold = any_bold or bold
        total_lines += _estimate_wrapped_lines(text, avail_width_pt, font_size, bold)

    if total_lines == 0 or max_font_size == 0:
        return

    needed_height_pt = total_lines * max_font_size * _LINE_HEIGHT_FACTOR
    if needed_height_pt <= avail_height_pt:
        return  # fits already, leave font sizes untouched

    scale = avail_height_pt / needed_height_pt
    scale = max(_MIN_SCALE, min(_MAX_SCALE, scale))
    # Round down to the nearest 5% -- matches PowerPoint's own step size and
    # keeps us slightly conservative rather than exactly on the boundary.
    scale = int(scale * 20) / 20.0

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            current = _run_font_size_pt(run, paragraph)
            new_size = max(_MIN_FONT_SIZE_PT, round(current * scale, 1))
            run.font.size = Pt(new_size)


def _clear_autofit_cache(text_frame):
    """
    Remove PowerPoint's stale cached autofit scale so nothing double-applies
    on top of the font sizes we just computed ourselves.
    """
    try:
        text_frame.word_wrap = True
    except Exception:
        pass

    body = text_frame._txBody
    bodyPr = body.find(f"{{{_A_NS}}}bodyPr")
    if bodyPr is None:
        return
    normAutofit = bodyPr.find(f"{{{_A_NS}}}normAutofit")
    if normAutofit is not None:
        for attr in ("fontScale", "lnSpcReduction"):
            normAutofit.attrib.pop(attr, None)


# ── replacement ──────────────────────────────────────────────────────────────
def _replace_in_text_frame(shape, text_frame, replacements, styling):
    changed = False

    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            continue

        # paragraph.text merges runs, so it catches placeholders that python-pptx
        # split across multiple runs (e.g. "{{Title}}" stored as "{{Ti" + "tle}}").
        para_text = paragraph.text
        combined = "".join(r.text for r in runs)

        present = [k for k in replacements if (k in combined or k in para_text)]
        if not present:
            continue

        # Prefer the merged-run string; fall back to paragraph.text if the
        # placeholder lives in XML python-pptx doesn't surface as runs.
        source = combined if any(k in combined for k in present) else para_text

        new_text = source
        for key in present:
            value = replacements[key]
            new_text = new_text.replace(key, "" if value is None else str(value))

        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""

        # Apply the style registered for the first placeholder in this paragraph.
        style = (styling or {}).get(present[0], {})
        for r in paragraph.runs:
            _apply_style(r, style)

        changed = True

    if changed:
        _clear_autofit_cache(text_frame)
        _shrink_to_fit(shape, text_frame)


def _fill_shape(shape, replacements, styling):
    # Groups (recurse into children)
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            _fill_shape(s, replacements, styling)

    # Tables
    if hasattr(shape, "table"):
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        _replace_in_text_frame(cell, cell.text_frame, replacements, styling)
        except Exception:
            pass

    # Plain text boxes / shapes with text
    if hasattr(shape, "text_frame"):
        try:
            _replace_in_text_frame(shape, shape.text_frame, replacements, styling)
        except Exception:
            pass


def fill_template(prs, replacements: dict, styling: dict = None):
    """Apply replacements + styling to every slide in an open Presentation."""
    styling = styling or {}
    replacements = replacements or {}
    for slide in prs.slides:
        for shape in slide.shapes:
            _fill_shape(shape, replacements, styling)
    return prs