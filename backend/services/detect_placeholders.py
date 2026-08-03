"""
Detect {{placeholder}} tokens in a template and return, for each one, a
normalized position/size box (as % of slide) plus basic font styling. The
frontend uses these boxes to lay its live-edit overlay over the rendered slide.

Supports both .pptx templates (via python-pptx) and .svg templates. The
download goes through the shared, guarded fetcher so this path enforces the
same https / size limits as preview and render.
"""

import os
import re
import tempfile
from xml.etree import ElementTree as ET

from pptx import Presentation

from services.template_fetch import download_template

_PLACEHOLDER_RE = re.compile(r"\{\{.*?\}\}")


# ── SVG templates ────────────────────────────────────────────────────────────
def extract_svg_placeholders(svg_content):
    """Extract {{placeholder}} patterns from SVG text elements."""
    placeholders_map = {}

    try:
        root = ET.fromstring(svg_content)

        # Slide dimensions (viewBox wins over width/height when present).
        width_attr = root.get("width", "1000")
        height_attr = root.get("height", "1000")
        svg_width = float(re.sub(r"[^\d.]", "", str(width_attr)) or 1000)
        svg_height = float(re.sub(r"[^\d.]", "", str(height_attr)) or 1000)

        viewBox = root.get("viewBox", "")
        if viewBox:
            parts = viewBox.strip().split()
            if len(parts) == 4:
                svg_width = float(parts[2])
                svg_height = float(parts[3])

        for elem in root.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            if tag not in ("text", "tspan", "flowRoot", "flowPara"):
                continue

            # Collect full text including tspan children.
            full_text = elem.text or ""
            for child in elem.iter():
                if child is not elem:
                    full_text += (child.text or "") + (child.tail or "")

            matches = _PLACEHOLDER_RE.findall(full_text)
            if not matches:
                continue

            x = float(elem.get("x", 0) or 0)
            y = float(elem.get("y", 0) or 0)

            # Fall back to a transform translate() when there's no direct x/y.
            t_match = re.search(
                r"translate\(([-\d.]+)[,\s]+([-\d.]+)\)", elem.get("transform", "")
            )
            if t_match:
                x += float(t_match.group(1))
                y += float(t_match.group(2))

            font_size = float(re.sub(r"[^\d.]", "", elem.get("font-size", "18")) or 18)
            font_color = elem.get("fill", "#1e293b")
            font_weight = elem.get("font-weight", "normal")
            font_style = elem.get("font-style", "normal")
            text_anchor = elem.get("text-anchor", "start")
            align_val = {"middle": "CENTER", "end": "RIGHT", "start": "LEFT"}.get(
                text_anchor, "LEFT"
            )

            estimated_width = max(len(full_text) * font_size * 0.6, font_size * 5)

            box = {
                "left": round((x / svg_width) * 100, 2),
                "top": round((y / svg_height) * 100, 2),
                "width": round(min((estimated_width / svg_width) * 100, 50), 2),
                "height": round((font_size * 1.5 / svg_height) * 100, 2),
                "fontSize": font_size,
                "align": align_val,
                "color": font_color if font_color.startswith("#") else "#1e293b",
                "isBold": font_weight in ("bold", "700", "800", "900"),
                "isItalic": font_style == "italic",
            }
            for match in matches:
                placeholders_map.setdefault(match, box)

    except ET.ParseError as e:
        raise ValueError(f"Invalid SVG file: {e}")

    return placeholders_map, svg_width, svg_height


# ── PPTX templates ───────────────────────────────────────────────────────────
def extract_placeholders(
    shape, slide_width, slide_height, placeholders_map,
    scale_x=1.0, scale_y=1.0, offset_x=0.0, offset_y=0.0,
):
    # Groups: recurse into children, tracking the child coordinate transform.
    if hasattr(shape, "shapes"):
        group_left_slide = (shape.left or 0) * scale_x + offset_x
        group_top_slide = (shape.top or 0) * scale_y + offset_y

        try:
            chOff = shape._element.xpath("./p:grpSpPr/a:xfrm/a:chOff")
            chOff_x = int(chOff[0].get("x")) if chOff else (shape.left or 0)
            chOff_y = int(chOff[0].get("y")) if chOff else (shape.top or 0)

            chExt = shape._element.xpath("./p:grpSpPr/a:xfrm/a:chExt")
            chExt_cx = int(chExt[0].get("cx")) if chExt else (shape.width or 1)
            chExt_cy = int(chExt[0].get("cy")) if chExt else (shape.height or 1)
        except Exception:
            chOff_x, chOff_y = (shape.left or 0), (shape.top or 0)
            chExt_cx, chExt_cy = (shape.width or 1), (shape.height or 1)

        new_scale_x = ((shape.width or 1) / (chExt_cx or 1)) * scale_x
        new_scale_y = ((shape.height or 1) / (chExt_cy or 1)) * scale_y
        new_offset_x = group_left_slide - chOff_x * new_scale_x
        new_offset_y = group_top_slide - chOff_y * new_scale_y

        for s in shape.shapes:
            extract_placeholders(
                s, slide_width, slide_height, placeholders_map,
                new_scale_x, new_scale_y, new_offset_x, new_offset_y,
            )

    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame"):
                    _extract_from_text_frame(
                        cell.text_frame, shape, slide_width, slide_height,
                        placeholders_map, scale_x, scale_y, offset_x, offset_y,
                    )

    if hasattr(shape, "text_frame"):
        _extract_from_text_frame(
            shape.text_frame, shape, slide_width, slide_height,
            placeholders_map, scale_x, scale_y, offset_x, offset_y,
        )


def _extract_from_text_frame(
    text_frame, shape, slide_width, slide_height, placeholders_map,
    scale_x, scale_y, offset_x, offset_y,
):
    for paragraph in text_frame.paragraphs:
        # Use paragraph.text (merged across all runs) to detect placeholders.
        # A single placeholder like {{Lorem Ipsum 2}} can be split across
        # multiple runs in the XML, so searching run-by-run misses them.
        matches = _PLACEHOLDER_RE.findall(paragraph.text)
        if not matches:
            continue

        left_val = (shape.left or 0) * scale_x + offset_x
        top_val = (shape.top or 0) * scale_y + offset_y
        width_val = (shape.width or 0) * scale_x
        height_val = (shape.height or 0) * scale_y

        left_pct = (left_val / slide_width) * 100
        top_pct = (top_val / slide_height) * 100
        width_pct = (width_val / slide_width) * 100 if width_val else 10
        height_pct = (height_val / slide_height) * 100 if height_val else 10

        font_size_pt = None
        font_color_hex = None
        is_bold = False
        is_italic = False

        # Scan all runs for font metadata -- even split runs carry valid styling.
        for run in paragraph.runs:
            if run.font:
                if not font_size_pt and run.font.size:
                    font_size_pt = run.font.size.pt
                if not font_color_hex and hasattr(run.font, "color"):
                    try:
                        if getattr(run.font.color, "rgb", None):
                            font_color_hex = f"#{run.font.color.rgb}"
                    except Exception:
                        pass
                if run.font.bold:
                    is_bold = True
                if run.font.italic:
                    is_italic = True

        if not font_size_pt and paragraph.font and paragraph.font.size:
            font_size_pt = paragraph.font.size.pt
        font_size_pt = font_size_pt or 18
        
        # Apply any scaling from parent groups to get the true visual size
        effective_font_size = font_size_pt * scale_y

        font_color_hex = font_color_hex or "#1e293b"

        align_val = _paragraph_alignment(paragraph)

        box = {
            "left": left_pct,
            "top": top_pct,
            "width": width_pct,
            "height": height_pct,
            "fontSize": round(effective_font_size, 1),
            "align": align_val,
            "color": font_color_hex,
            "isBold": is_bold,
            "isItalic": is_italic,
        }
        for match in matches:
            placeholders_map.setdefault(match, box)


def _paragraph_alignment(paragraph):
    raw_align = paragraph.alignment
    if raw_align is None:
        try:
            pPr = paragraph._p.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
            )
            algn = pPr.get("algn") if pPr is not None else None
            return {"r": "RIGHT", "l": "LEFT"}.get(algn, "CENTER")
        except Exception:
            return "CENTER"
    align_str = str(raw_align)
    if "CENTER" in align_str:
        return "CENTER"
    if "RIGHT" in align_str:
        return "RIGHT"
    return "LEFT"


# ── entry point ──────────────────────────────────────────────────────────────
def detect_placeholders(template_url: str):
    work_dir = tempfile.mkdtemp(prefix="detect_")
    tmp_path = os.path.join(work_dir, "template.bin")
    try:
        download_template(template_url, tmp_path)

        with open(tmp_path, "rb") as f:
            head = f.read(512)

        url_lower = template_url.lower().split("?")[0]
        is_svg = url_lower.endswith(".svg") or b"<svg" in head

        if is_svg:
            with open(tmp_path, "rb") as f:
                svg_text = f.read().decode("utf-8", errors="replace")
            placeholders_map, svg_width, svg_height = extract_svg_placeholders(svg_text)
            return {
                "success": True,
                "file_type": "svg",
                "placeholders": list(placeholders_map.keys()),
                "placeholder_mappings": placeholders_map,
                "slide_width": svg_width,
                "slide_height": svg_height,
            }

        prs = Presentation(tmp_path)
        placeholders_map = {}
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides):
            before = set(placeholders_map.keys())
            for shape in slide.shapes:
                extract_placeholders(shape, slide_width, slide_height, placeholders_map)
            for ph in placeholders_map:
                if ph not in before and "slide_index" not in placeholders_map[ph]:
                    placeholders_map[ph]["slide_index"] = slide_idx

        return {
            "success": True,
            "file_type": "pptx",
            "placeholders": list(placeholders_map.keys()),
            "placeholder_mappings": placeholders_map,
            "slide_width": slide_width,
            "slide_height": slide_height,
        }
    finally:
        import shutil
        shutil.rmtree(work_dir, ignore_errors=True)